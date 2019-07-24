from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from enum import Enum

class Text_Type ( Enum ):
	Title = 1
	Body = 2
	Text_Box = 3
	Shape = 4

# The Return_Text function parses a pptx file into a 2-dimensional array.
# The first dimension indexes into the slide deck (ie index i is slide i + 1).
# The second dimension corresponds to an array of text elements that are
# present on the slide determined by the first coordinate. The format for each
# element will be ( Type, Text ) where Type is an enum that takes on either
# Title(1), Body(2), Text_Box(3), or Shape(4). The Text is simply a string of
# the contents of contained in that particular text element.
# param file_name : the file path to the pptx file to parse.
# returns : the 2-dimensional [slide][text_element] array described above.
def Return_Text( file_name ):
	prs = Presentation( file_name )
	slide_list = []

	for slide in prs.slides:
		elements = []
		for shape in slide.shapes:
			if shape.has_text_frame:
				for para in shape.text_frame.paragraphs:
					# Don't do anything with empty paragraphs.
					if para.text == "":
						continue
					elif shape.is_placeholder: # Shape is a placeholder. Title vs Body.
						if shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
							elements.append( (Text_Type.Title, para.text) )
						elif shape.placeholder_format.type == PP_PLACEHOLDER.OBJECT:
							elements.append( (Text_Type.Body, para.text) )
					else: # Shape is not a placeholder. It's a generic shape. Textbox vs Shape.
						if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
							elements.append( ( Text_Type.Text_Box, para.text) )
						elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
							elements.append( ( Text_Type.Shape, para.text) )
		# Append if elements were found on slide.
		if len(elements) != 0:
			slide_list.append(elements)
	return slide_list
           
# for slide in prs.slides:
#     for shape in slide.shapes:
#         if 'Picture' in shape.name:
#             picture = shape
#             image = picture.image
#             image_file_bytes = image.blob
#             file_extension = image.ext
#             file_path=image.filename
#             print(file_path)


# for shape in slide.shapes:
#     if shape.shape_type != MSO_SHAPE_TYPE:
#         continue
#     picture = shape
#     print(picture._pic.nvPicPr.cNvPr.get('descr'))
